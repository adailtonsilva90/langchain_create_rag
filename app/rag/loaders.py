import os
import json
import base64
from typing import List, Dict, Any
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.messages import HumanMessage

def get_llm():
    model_name = os.getenv("GEMINI_ANALYSIS_MODEL", "models/gemini-2.5-flash")
    return ChatGoogleGenerativeAI(model=model_name, temperature=0)

def extract_raw_text(file_bytes: bytes, filename: str) -> str:
    """
    Extracts ONLY the raw text from a document using Gemini.
    """
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == ".pptx":
        try:
            import io
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text)
            return "\n".join(text_content)
        except Exception as e:
            print(f"Error extracting PPTX raw text: {e}")
            return f"Error: {e}"

    llm = get_llm()
    mime_type = "application/pdf" if ext == ".pdf" else f"image/{ext[1:]}"
    if ext in [".jpg", ".jpeg"]: mime_type = "image/jpeg"
    
    data_b64 = base64.b64encode(file_bytes).decode("utf-8")
    
    message = HumanMessage(
        content=[
            {"type": "text", "text": "Extract all the text from this document literally and completely. Return only the extracted text."},
            {"type": "media", "mime_type": mime_type, "data": data_b64}
        ]
    )

    try:
        response = llm.invoke([message])
        return response.content.strip()
    except Exception as e:
        print(f"Error in extract_raw_text: {e}")
        return f"Error: {e}"

def extract_semantic_content(file_bytes: bytes, filename: str) -> Dict[str, Any]:
    """
    Extracts text and semantic sections (JSON) using Gemini.
    Used for the ingestion pipeline.
    """
    llm = get_llm()
    ext = os.path.splitext(filename)[1].lower()
    
    prompt = (
        "Analyze the document and extract the text in a structured way.\n"
        "Identify natural semantic sections (e.g. Experience, Education, Chapters).\n\n"
        "Return EXCLUSIVELY a JSON in the format:\n"
        "{\n"
        "  \"raw_text\": \"full text\",\n"
        "  \"sections\": [\n"
        "    {\"category\": \"name\", \"content\": \"content\"}\n"
        "  ]\n"
        "}"
    )

    if ext == ".pptx":
        try:
            import io
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())
            
            extracted_text = "\n".join(text_content)
            
            message = HumanMessage(
                content=[
                    {"type": "text", "text": f"{prompt}\n\nHere is the text extracted from the presentation:\n\n{extracted_text}"}
                ]
            )
        except Exception as e:
            print(f"Error reading PPTX offline: {e}")
            return {"raw_text": "", "sections": []}
    else:
        mime_type = "application/pdf" if ext == ".pdf" else f"image/{ext[1:]}"
        if ext in [".jpg", ".jpeg"]: mime_type = "image/jpeg"
        
        data_b64 = base64.b64encode(file_bytes).decode("utf-8")
        
        message = HumanMessage(
            content=[
                {"type": "text", "text": prompt},
                {"type": "media", "mime_type": mime_type, "data": data_b64}
            ]
        )

    try:
        response = llm.invoke([message])
        content = response.content.strip()
        if content.startswith("```json"):
            content = content.replace("```json", "").replace("```", "").strip()
        elif content.startswith("```"):
            content = content.replace("```", "").strip()
        
        return json.loads(content)
    except Exception as e:
        print(f"Error in extract_semantic_content: {e}")
        return {"raw_text": "", "sections": []}
